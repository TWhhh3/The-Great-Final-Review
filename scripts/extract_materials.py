from pathlib import Path
import json


ROOT = Path(__file__).resolve().parents[1]
MATERIALS_DIR = ROOT / "materials"
EXTRACTED_DIR = ROOT / "extracted"

SUPPORTED_TYPES = {".pdf", ".pptx", ".docx", ".txt", ".md", ".png", ".jpg", ".jpeg"}

MANUAL_IMAGE_TEXT = {
    "模拟试卷（题型说明）/第一页.png": """上届考题（部分）示例

一、选择题（每题 2 分，共 30 分）（15 小题）

1、必须在中断服务程序中，用软件清除的中断请求标志是（ ）。
（A）TI、TF1    （B）RI、IE1    （C）TF0、TF1    （D）TI、RI

2、在 MCS-51 单片机的并行 I/O 口中，可分时作地址线和数据线的是（ ）。
（A）P0 口    （B）P1 口    （C）P2 口    （D）P3 口

3、若 CPU 采用 13 根地址线寻址时，其寻址范围是（ ）。
（A）1KB    （B）2KB    （C）4KB    （D）8KB

4、MCS-51 单片机进行系统扩展时，设置了一个锁存器 74LS373 用于锁存（ ）。
（A）输入/输出数据    （B）低 8 位地址
（C）高 8 位地址      （D）控制信号

5、在 MCS-51 单片机中，当晶振频率采用 2MHZ 时，一个机器周期等于（ ）微秒。
（A）3    （B）4    （C）5    （D）6

……

二、填空题（每空 1 分，共 10 分）

1、十进制数 0、-128、-1 的补码分别为__________，__________，__________。

2、51 单片机 CPU 对片外 RAM 的访问只能采用__________寻址方式，并且片外 RAM 的数据只能与__________之间进行数据传送。

……

三、程序阅读题（12 分）（共两小题）

1、阅读下列程序，回答问题

ORG 2000H
START: MOV R0, #05H
       MOV R1, #30H ; 30H 为一组带符号数的存储起始地址。
       MOV R2, #00H
LOOP:  MOV A, @R1
       JNB ACC.7, LOOP1
       INC R2
LOOP1: INC R1
       DJNZ R0, LOOP
       SJMP $
       END

（1）程序的功能是________________________________。
（2）若从 30H 至 34H 单元的内容依次为 A8H、7BH、96H、03H、45H，则程序执行后 R2=________。

……

四、简答题（10 分）（两小题）

1、子程序或中断服务程序中，为什么要保护现场？如何保护和恢复现场？（5 分）
""",
    "模拟试卷（题型说明）/第二页.png": """五、系统扩展（8 分）（1 题）

如图，给 8051 单片机扩展一片 EPROM 2732（4KB）作为程序存储器、一片 SRAM 6232（4KB）作为片外数据存储器。

图中主要器件和信号：
- 8051
- 74LS373 锁存器
- 2732 程序存储器
- 6232 数据存储器
- 地址线、数据线、ALE、PSEN、RD、WR、EA 等连接关系见原图

（1）完成图中未完成的连接线。

（2）确定程序存储器 2732 的地址范围：____________________；以及数据存储器 6232 的地址范围____________________。
（注：将无关的地址位统一取 0）

（3）编程将程序存储器 2732 最低地址单元的内容传送到数据存储器 6232 最高地址的单元中。

六、汇编语言程序设计（30 分）（3～4 小题）

（26 年程序设计题：汇编语言是基础，但大部分都可选汇编或 C51，不强制用 C51）

如图，8051 单片机使用 74LS245 扩展了一个输入口（外设 A 口）。

1）确定外设 A 口的地址。

2）编一个程序，使用软件延时的方法，每 20ms 读取外设 A 口的值，并采用工作方式 2 将读到的值从串口发送出去。
（提示：串口通讯采用波特率固定的工作方式 2 时，不需计算波特率及设置定时器 T1 的初值。20ms 延时子程序不用写出）

图中主要器件和信号：
- 8051
- 74LS245
- 外设 A 口（8 位）
- P2.6、P2.5、ALE、P0、WR、RD、EA 等连接关系见原图
""",
}


def read_text_file(path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "gbk"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    print(f"[无法处理] 文本编码无法识别：{path}")
    return ""


def extract_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        print("[缺少依赖] 请先安装 pypdf：pip install -r requirements.txt")
        return ""

    parts = []
    try:
        reader = PdfReader(str(path))
        for page_number, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            parts.append(f"=== PDF 第 {page_number} 页 ===")
            if text:
                parts.append(text)
            else:
                parts.append("资料未明确说明：本页未提取到可直接读取的文字，可能是扫描页，需要 OCR。")
            parts.append("")
    except Exception as error:
        print(f"[无法处理] PDF 提取失败：{path}，原因：{error}")

    return "\n".join(parts)


def extract_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        print("[缺少依赖] 请先安装 python-docx：pip install -r requirements.txt")
        return ""

    parts = []
    try:
        document = Document(str(path))

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if text:
                parts.append(text)

        for table_number, table in enumerate(document.tables, start=1):
            parts.append("")
            parts.append(f"=== Word 表格 {table_number} ===")
            for row in table.rows:
                cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                line = " | ".join(cell for cell in cells if cell)
                if line:
                    parts.append(line)
    except Exception as error:
        print(f"[无法处理] DOCX 提取失败：{path}，原因：{error}")

    return "\n".join(parts)


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        print("[缺少依赖] 请先安装 python-pptx：pip install -r requirements.txt")
        return ""

    parts = []
    try:
        presentation = Presentation(str(path))
        for slide_number, slide in enumerate(presentation.slides, start=1):
            title = ""
            if slide.shapes.title is not None:
                title = slide.shapes.title.text.strip()

            header = f"=== PPT 第 {slide_number} 张幻灯片"
            if title:
                header += f"：{title}"
            header += " ==="
            parts.append(header)

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        parts.append(text)
            parts.append("")
    except Exception as error:
        print(f"[无法处理] PPTX 提取失败：{path}，原因：{error}")

    return "\n".join(parts)


def extract_image(path: Path) -> str:
    relative = path.relative_to(MATERIALS_DIR).as_posix()
    if relative in MANUAL_IMAGE_TEXT:
        print(f"[图片转写] 已根据图片内容转写：{path.relative_to(ROOT).as_posix()}")
        return "说明：本文件为 PNG 图片，以下文字根据图片内容转写。\n\n" + MANUAL_IMAGE_TEXT[relative]

    print(f"[无法处理] 图片需要 OCR 或人工转写：{path.relative_to(ROOT).as_posix()}")
    return "资料未明确说明：该图片尚未配置 OCR 或人工转写内容。"


def output_path_for(source_path: Path) -> Path:
    relative = source_path.relative_to(MATERIALS_DIR)
    return (EXTRACTED_DIR / relative).with_suffix(".txt")


def extract_one(path: Path) -> dict:
    extension = path.suffix.lower()
    relative_source = path.relative_to(ROOT).as_posix()

    if extension not in SUPPORTED_TYPES:
        print(f"[跳过] 暂不支持此文件类型：{relative_source}")
        return {
            "source": relative_source,
            "output": "",
            "status": "skipped",
            "message": f"unsupported type: {extension}",
        }

    if extension == ".pdf":
        text = extract_pdf(path)
    elif extension == ".docx":
        text = extract_docx(path)
    elif extension == ".pptx":
        text = extract_pptx(path)
    elif extension in {".png", ".jpg", ".jpeg"}:
        text = extract_image(path)
    else:
        text = read_text_file(path)

    out_path = output_path_for(path)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    header = [
        f"来源文件：{relative_source}",
        f"文件类型：{extension}",
        "",
    ]
    final_text = "\n".join(header) + (text if text else "资料未明确说明：没有提取到文字。")
    out_path.write_text(final_text, encoding="utf-8")

    print(f"[完成] {relative_source} -> {out_path.relative_to(ROOT).as_posix()}")
    return {
        "source": relative_source,
        "output": out_path.relative_to(ROOT).as_posix(),
        "status": "extracted",
        "chars": len(text),
    }


def main() -> None:
    if not MATERIALS_DIR.exists():
        print("[错误] 找不到 materials 文件夹。")
        return

    EXTRACTED_DIR.mkdir(parents=True, exist_ok=True)

    results = []
    for path in sorted(MATERIALS_DIR.rglob("*")):
        if path.is_file():
            results.append(extract_one(path))

    index_path = EXTRACTED_DIR / "source_index.json"
    index_path.write_text(json.dumps(results, ensure_ascii=False, indent=2), encoding="utf-8")

    extracted_count = sum(1 for item in results if item["status"] == "extracted")
    skipped_count = sum(1 for item in results if item["status"] == "skipped")
    print("")
    print(f"提取完成：成功 {extracted_count} 个，跳过 {skipped_count} 个。")
    print(f"索引文件：{index_path.relative_to(ROOT).as_posix()}")


if __name__ == "__main__":
    main()
